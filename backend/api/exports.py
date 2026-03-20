from datetime import datetime
import json
from pathlib import Path
from typing import List

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from config import settings
from database import get_db
from models.content import ContentItem, ContentVariant
from services.storage import storage_service


router = APIRouter(prefix="/exports", tags=["exports"])


class ExportRequest(BaseModel):
    content_item_ids: List[str] = Field(default_factory=list)
    format: str = Field(default="md")  # md | html | pdf | pptx | json


@router.post("")
async def export_content(request: ExportRequest, db: AsyncSession = Depends(get_db)):
    if not request.content_item_ids:
        raise HTTPException(status_code=400, detail="content_item_ids 不能为空")

    result = await db.execute(
        select(ContentItem).where(ContentItem.id.in_(request.content_item_ids))
    )
    items = list(result.scalars().all())
    if not items:
        raise HTTPException(status_code=404, detail="未找到可导出的内容条目")

    bundle = []
    for item in items:
        variants_result = await db.execute(
            select(ContentVariant).where(ContentVariant.content_item_id == item.id)
        )
        variants = list(variants_result.scalars().all())
        bundle.append({"item": item, "variants": variants})

    fmt = request.format.lower()
    if fmt == "md":
        payload = _to_markdown(bundle).encode("utf-8")
        file_name = f"geo_export_{_ts()}.md"
        media_type = "text/markdown"
    elif fmt == "html":
        payload = _to_html(bundle).encode("utf-8")
        file_name = f"geo_export_{_ts()}.html"
        media_type = "text/html"
    elif fmt == "json":
        payload = _to_json(bundle).encode("utf-8")
        file_name = f"geo_export_{_ts()}.json"
        media_type = "application/json"
    elif fmt == "pdf":
        payload = _to_pdf(bundle)
        file_name = f"geo_export_{_ts()}.pdf"
        media_type = "application/pdf"
    elif fmt == "pptx":
        payload = _to_pptx(bundle)
        file_name = f"geo_export_{_ts()}.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        raise HTTPException(status_code=400, detail="不支持的导出格式")

    stored_path = storage_service.save_export_bytes(file_name, payload)
    if stored_path.startswith("s3://"):
        raise HTTPException(
            status_code=501,
            detail="当前构建未启用 S3 导出下载代理。",
        )
    return FileResponse(path=stored_path, media_type=media_type, filename=file_name)


def _ts() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def _to_markdown(bundle: list) -> str:
    out = [f"# GEO 导出文件\n\n生成时间：{datetime.now().isoformat()}\n"]
    for row in bundle:
        item = row["item"]
        out.append(f"\n## 主题：{item.topic}\n")
        out.append(f"- 内容类型：{item.content_type}\n- 状态：{item.status}\n")
        for v in row["variants"]:
            out.append(f"\n### [{v.platform}] {v.title}\n\n{v.body}\n")
    return "\n".join(out)


def _to_html(bundle: list) -> str:
    html = [
        "<!doctype html><html><head><meta charset='utf-8'><title>GEO 导出文件</title></head><body>",
        f"<h1>GEO 导出文件</h1><p>生成时间：{datetime.now().isoformat()}</p>",
    ]
    for row in bundle:
        item = row["item"]
        html.append(f"<h2>主题：{item.topic}</h2><p>类型：{item.content_type} | 状态：{item.status}</p>")
        for v in row["variants"]:
            body = (v.body or "").replace("\n", "<br>")
            html.append(f"<h3>[{v.platform}] {v.title}</h3><div>{body}</div>")
    html.append("</body></html>")
    return "\n".join(html)


def _to_json(bundle: list) -> str:
    payload = []
    for row in bundle:
        item = row["item"]
        payload.append(
            {
                "item_id": item.id,
                "topic": item.topic,
                "content_type": item.content_type,
                "status": item.status,
                "variants": [
                    {
                        "id": v.id,
                        "platform": v.platform,
                        "title": v.title,
                        "body": v.body,
                        "tags": v.tags,
                    }
                    for v in row["variants"]
                ],
            }
        )
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _to_pdf(bundle: list) -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.pdfgen import canvas
    except Exception:
        return _to_markdown(bundle).encode("utf-8")

    export_path = Path(settings.EXPORT_DIR) / f"tmp_{_ts()}.pdf"
    c = canvas.Canvas(str(export_path), pagesize=A4)
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    c.setFont("STSong-Light", 11)
    y = 800
    c.drawString(40, y, "GEO 导出文件")
    y -= 24

    for row in bundle:
        item = row["item"]
        lines = [f"主题: {item.topic}", f"类型: {item.content_type}  状态: {item.status}"]
        for v in row["variants"]:
            lines.append(f"[{v.platform}] {v.title}")
            lines.extend((v.body or "")[:1000].splitlines()[:10])
            lines.append("")
        for line in lines:
            if y < 80:
                c.showPage()
                c.setFont("STSong-Light", 11)
                y = 800
            c.drawString(40, y, line[:90])
            y -= 16
        y -= 8
    c.save()
    data = export_path.read_bytes()
    export_path.unlink(missing_ok=True)
    return data


def _to_pptx(bundle: list) -> bytes:
    try:
        from pptx import Presentation
    except Exception:
        return _to_markdown(bundle).encode("utf-8")

    export_path = Path(settings.EXPORT_DIR) / f"tmp_{_ts()}.pptx"
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "GEO 导出文件"
    title_slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    for row in bundle:
        item = row["item"]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = item.topic or "未命名主题"
        body = slide.shapes.placeholders[1].text_frame
        body.text = f"类型: {item.content_type} | 状态: {item.status}"
        for v in row["variants"]:
            p = body.add_paragraph()
            p.text = f"[{v.platform}] {v.title}"
            p.level = 0
            p2 = body.add_paragraph()
            p2.text = (v.summary or v.body[:120]).replace("\n", " ")
            p2.level = 1
    prs.save(str(export_path))
    data = export_path.read_bytes()
    export_path.unlink(missing_ok=True)
    return data
